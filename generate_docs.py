import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DocxPt
import os

def create_presentation():
    prs = Presentation()
    
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Conception et Simulation d'un Drone Parapluie Autonome"
    subtitle.text = "Système de suivi intelligent basé sur la vision par ordinateur et l'IA\n\nPrésenté par : Firas Mrabet\nSimulation sous Webots R2025a"
    
    def add_slide(title_str, bullet_points):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title_str
        tf = s.placeholders[1].text_frame
        for i, point in enumerate(bullet_points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            if isinstance(point, list):
                p.text = point[0]
                p.level = 1
            else:
                p.text = point
                p.level = 0
    
    add_slide("Contexte & Problématique", [
        "Le Besoin: Protection météorologique (pluie, soleil) 100% mains-libres.",
        "Cibles visées: Piétons, travailleurs en extérieur, personnes à mobilité réduite.",
        "Problématique:",
        ["Trouver, identifier et suivre sa cible sans collision.",
         "Maintenir une position de couverture optimale (au-dessus de la tête)."],
        "Notre Solution: Un modèle robotique simulé sous Webots, combinant le contrôle PID et l'intelligence artificielle visuelle (YOLO/OpenCV)."
    ])
    
    add_slide("Objectifs du Projet", [
        "Concevoir un environnement de simulation réaliste (Webots).",
        "Développer un pipeline de vision hybride rapide et robuste.",
        "Créer une Machine à États (FSM) pour automatiser le cycle : Décollage -> Recherche -> Approche -> Suivi.",
        "Surmonter les contraintes physiques des capteurs internes du modèle 3D."
    ])
    
    add_slide("Technologies et Outils Utilisés", [
        "Environnement 3D : Cyberbotics Webots (Physique, Capteurs).",
        "Programmation : Python 3 (API Webots orientée objet).",
        "Intelligence Artificielle & Traitement d'Image :",
        ["Ultralytics YOLOv8 : Réseau de neurones profond (Deep Learning) pour la classification d'humains.",
         "OpenCV (cv2) : Algorithmes ultra-rapides et masques colorimétriques (Espace HSV)."],
        "Gestion de Version : Git & GitHub (CI/CD Collaboratif)."
    ])
    
    add_slide("Architecture du Drone (Mavic 2 Pro)", [
        "Moteurs principaux : 4 rotors (Poussée, Roulis, Tangage, Lacet).",
        "Système de Caméra Double (Solution innovante aux butées mécaniques) :",
        ["1. Caméra Frontale sur Nacelle motorisée : Cherche au loin (inclinée à -17° pour balayage 360).",
         "2. Caméra Ventrale Fixe (Track_Camera) : Fixée sous le châssis à -90°, essentielle pour le suivi zénithal de précision."],
        "Capteurs de Vol Régulateurs : Gyroscope, IMU, GPS."
    ])
    
    add_slide("L'Intelligence Artificielle (Architecture)", [
        "Problème: L'IA neuronale (YOLOv8) est trop coûteuse en ressources CPU dans le simulateur s'il tourne à chaque ms.",
        "Solution - Méthode Optimisée 'Hybride Asynchrone' :",
        ["A) Détection Lointaine Active : YOLOv8 s'active lors du radar 360° pour confirmer l'humain.",
         "B) Suivi Haute Fréquence (PID) : Pour asservir les moteurs pendant la descente et la marche, on bascule sur HSV OpenCV."],
        "Résultat final : Le taux de framerate chute de 0%, l'environnement reste fluide."
    ])
    
    add_slide("Séquence de Vol Autonome", [
        "PHASE 1 (Décollage) : Montée en altitude (2.8m).",
        "PHASE 2 (Recherche) : Vol stationnaire, Rotation du drone à 360° sur lui-même (Radar).",
        "PHASE 3 (Approche) : Cible détectée (YOLO). Pivot face au piéton, vol de translation lent (-40cm).",
        "PHASE 4 (Descente) : Bascule sur flux vidéo inférieur. Descente douce sur cible.",
        "PHASE 5 (Suivi Synchronisé) : Ordre de marche IPC, suivi de la position au centre pixels."
    ])
    
    add_slide("Discussion & Perspectives", [
        "Bilan :",
        ["L'orchestration asynchrone YOLO+Couleurs combinée au 'Dual Camera' résout élégamment les limites informatiques et la modélisation des limites mécaniques du hardware."],
        "Perspectives futures :",
        ["Implémentation directe sur Nvidia Jetson ou un Raspberry Pi.",
         "Capteurs LiDAR pour évitement d'obstacles dynamiques.",
         "Intelligence Artificielle Prédictive (Kalman) pour gérer les accélérations/freinages du piéton."]
    ])
    
    prs.save("/home/firas.mrabet/Desktop/drone_parapluie_projet/docs/Presentation_Drone_Parapluie.pptx")

def create_report():
    doc = Document()
    doc.add_heading('Rapport de Projet : Simulation d\'un Drone Parapluie Autonome', 0)
    
    doc.add_heading('1. Introduction Générale', level=1)
    doc.add_paragraph("L'évolution rapide de la robotique de service ouvre la voie à de nouvelles interactions entre l'humain et les drones autonomes. Ce projet s'inscrit dans cette dynamique en proposant la conception et la simulation d'un \"Drone Parapluie\". L'objectif est de développer un système robotique capable de décoller, de détecter une personne cible de manière autonome, puis de se positionner au-dessus d'elle pour la protéger des intempéries (pluie, soleil) tout en se synchronisant avec sa marche.")
    doc.add_paragraph("La réalisation de ce système complexe nécessite une synergie entre le contrôle de vol classique et l'Intelligence Artificielle embarquée, le tout validé dans un environnement de simulation tridimensionnel (Cyberbotics Webots).")
    
    doc.add_heading('2. État de l\'Art et Contexte Technologique', level=1)
    doc.add_paragraph("Le suivi de cible par drone (Object Tracking) repose généralement sur des signaux de radiolocalisation (capteurs embarqués) ou sur la vision par ordinateur. Notre projet privilégie une approche 100% visuelle, libérant l'utilisateur de toute contrainte matérielle sur lui-même.")
    doc.add_paragraph("Pour ce faire, nous avons évalué et couplé les algorithmes de vision :")
    doc.add_paragraph("Deep Learning (YOLOv8) : L'état de l'art, extrêmement précis pour identifier la classe 'Person', mais très exigeant en puissance de calcul (CPU/GPU).", style='List Bullet')
    doc.add_paragraph("Filtres Colorimétriques (HSV - OpenCV) : Traitement mathématique d'espace vectoriel ultra-rapide (2 millisecondes), robuste à la latence, mais fragile si utilisé seul en open-world.", style='List Bullet')
    doc.add_paragraph("Afin de combiner justesse sémantique et vélocité pour le feedback PID du robot, ce projet propose une architecture logicielle hybride.")
    
    doc.add_heading('3. Architecture Robotique Multimédia (Simulation)', level=1)
    doc.add_heading('3.1. Le Drone et ses Actuateurs', level=2)
    doc.add_paragraph("Basé sur le modèle virtuel du constructeur DJI Mavic 2 Pro. Il incorpore quatre hélices motorisées, une unité centrale inertielle (IMU), des capteurs gyroscopiques, et un récepteur GPS de position absolue pour stabiliser son vol à haute fréquence.")
    
    doc.add_heading('3.2. Le Pipeline ' + "Dual-Camera" + ' (Innovation mécanique)', level=2)
    doc.add_paragraph("Pendant la modélisation R&D, une contrainte biomécanique majeure fut rencontrée : la nacelle stabilisatrice (motor gimbal) du Mavic est empêchée physiquement de pivoter sous la barre des -28 degrés vers l'axe -Z local. Il ne pouvait donc pas voir le sommet exact du crâne humain sans basculer physiquement l'appareil.")
    doc.add_paragraph("Notre architecture remédie à cela grâce au déploiement asynchrone de deux flux vidéo :")
    doc.add_paragraph("La Caméra de Base (Search) : Caméra du gimbal frontal pour scruter l'horizon (Phase d'approche).", style='List Bullet')
    doc.add_paragraph("La Caméra Zénithale Ventrale (Track) : Matériel ajouté spécifiquement et boulonné sous la canopée du parapluie, braqué fixement vers le sol (-90°) pur pour la phase d'alignement fin de couverture pluie.", style='List Bullet')
    
    doc.add_heading('4. Modélisation de la Machine à États (Control Flow)', level=1)
    doc.add_paragraph("L'intelligence orchestre la FSM (Finite State Machine) selon 5 états programmés :")
    doc.add_paragraph("PHASE 1 (Décollage) : Correction d'erreur altimètre PID stricte vers 2.8m.", style='List Number')
    doc.add_paragraph("PHASE 2 (Stationnaire Radar) : Maintien des coordonnées GPS (Hover), Balayage à 360° grâce à la rotation sur le lacet (Yaw vector) pour éveiller l'IA YOLOv8.", style='List Number')
    doc.add_paragraph("PHASE 3 (Capture de Cible) : Fixation de la source. Le composant s'oriente face au piéton, calcule son vecteur normatif, et entame une translation spatiale douce (alpha tracking).", style='List Number')
    doc.add_paragraph("PHASE 4 (Lâcher des balasts) : Basculement instantané des ressources logicielles GPU/CPU sur la caméra ventrale et descente graduelle au zénith (2.0m d'élévation).", style='List Number')
    doc.add_paragraph("PHASE 5 (PID Autotracking) : La boucle PID régulatrice tourne à plus de 32 Hz pour aligner continuellement les erreurs XY de la matrice image (barycentriques) par rapport aux forces propulsives.", style='List Number')
    
    doc.add_heading('5. Conclusion Analytique', level=1)
    doc.add_paragraph("L'accomplissement global de la chaîne FSM/CV valide profondément la pertinence du simulateur Webots pour modéliser des réseaux complexes. L'approche dual-camera et le fractionnement calculatoire YOLO/OpenCV démontrent au jury une réelle résolution de problème industriel asymétrique face aux contraintes temporelles strictes qu'attendent les nano-micro-ordinateurs d'edge-computing aéronautiques modernes.")
    
    doc.save("/home/firas.mrabet/Desktop/drone_parapluie_projet/docs/Rapport_Drone_Parapluie.docx")

if __name__ == "__main__":
    create_presentation()
    create_report()
    print("Files created successfully.")
